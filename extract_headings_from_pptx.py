from pptx import Presentation


def extract_slide_headings(pptx_file_path):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file_path)

    # Initialize a list to store the extracted headings
    headings = []

    # Iterate through slides and extract headings from slide layouts
    for slide in presentation.slides:
        if slide.shapes.title is not None:
            heading = slide.shapes.title.text

            # Replace non-breaking space (nbsp) and vertical tab (vt) with regular spaces
            heading = heading.replace('\xa0', ' ').replace('\x0b', ' ')

            if heading.strip():  # Check if the heading is not empty
                headings.append(heading)

    return headings


# Replace 'your_presentation.pptx' with the name of your PowerPoint file
pptx_file_path = '..\\files\\your_presentation.pptx'

# Extract slide headings from the PowerPoint file
headings = extract_slide_headings(pptx_file_path)

# Define the output file name
output_file = '..\\files\\extracted_headings.txt'

# Write the extracted headings to a text file with a new line after each heading
with open(output_file, 'w', encoding='utf-8') as file:
    for heading in enumerate(headings, start=1):
        file.write(f"{heading}\n\n")  # Add a new line after each heading

# Print a message indicating where the headings have been saved
print(f"Extracted slide headings saved to '{output_file}'")